from __future__ import annotations

"""
Chat file upload endpoint — extract text from uploaded files for AI context.

Files uploaded here are processed in-memory and never stored to S3 or the
database. This is intentional: chat uploads are ephemeral AI context, not
durable user documents.

Supported formats:
  Text:   .txt  .md  .html  .htm  .csv  .json  .xml  .py  .js  .ts
  Office: .docx  .xlsx  .xls  .pptx  .pdf
  Image:  .png  .jpg  .jpeg  .heic  .webp  .gif  (returns placeholder)
  Media:  .mp3  .mp4  .mov  .wav  .m4a       (returns placeholder)
  Zip:    .zip  (extracts and reads text files inside)
"""

import io
import logging
import zipfile
from pathlib import Path

from fastapi import APIRouter, Request, UploadFile, File, HTTPException
from database import get_db
from middleware.auth import get_current_user
from services.audit_service import write_audit
from services.ocr_service import extract_text as extract_image_text
from services.ocr_service import sniff_image_type
from services.vision_service import describe_image as describe_image_bytes
from tenant import get_school_id

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/chat", tags=["chat-upload"])

MAX_FILE_SIZE_BYTES = 20 * 1024 * 1024  # 20 MB
MAX_TEXT_LENGTH = 40_000                # chars forwarded to LLM context
# R9.4 (X8): zip-bomb guards — a single member and the total decompressed size
# are both capped using the archive DIRECTORY's declared sizes, BEFORE any member
# is actually decompressed.
MAX_ZIP_MEMBER_BYTES = 5 * 1024 * 1024      # 5 MB uncompressed per member
MAX_ZIP_TOTAL_BYTES = 50 * 1024 * 1024      # 50 MB uncompressed total
_READ_CHUNK_BYTES = 1024 * 1024             # 1 MB streaming read chunk
BLOCKED_EXTENSIONS = {
    ".exe", ".bat", ".sh", ".cmd", ".com", ".ps1",
    ".vbs", ".jar", ".msi", ".dll", ".bin", ".scr",
}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".yaml", ".yml",
    ".sql", ".log",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".heic", ".webp", ".gif", ".bmp", ".tiff"}
MEDIA_EXTENSIONS = {".mp3", ".mp4", ".mov", ".wav", ".m4a", ".aac", ".avi", ".mkv"}


def _read_plain_text(data: bytes, filename: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes, filename: str) -> str:
    try:
        import pypdf  # noqa: F401 — optional dep
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text.strip())
        if pages:
            return "\n\n".join(pages)
        return f"[PDF: {filename} — no extractable text found (may be scanned)]"
    except ImportError:
        logger.warning("pypdf not installed; cannot extract PDF text")
        return f"[PDF: {filename} — install pypdf to enable text extraction]"
    except Exception as exc:
        logger.warning("PDF extraction error for %s: %s", filename, exc)
        return f"[PDF: {filename} — extraction error: {exc}]"


def _extract_docx(data: bytes, filename: str) -> str:
    try:
        from docx import Document  # noqa: F401
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract table cells
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs) if paragraphs else f"[DOCX: {filename} — no text found]"
    except ImportError:
        logger.warning("python-docx not installed; cannot extract .docx text")
        return f"[DOCX: {filename} — install python-docx to enable text extraction]"
    except Exception as exc:
        logger.warning("DOCX extraction error for %s: %s", filename, exc)
        return f"[DOCX: {filename} — extraction error: {exc}]"


def _extract_xlsx(data: bytes, filename: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in row_vals):
                    rows.append(" | ".join(row_vals))
            if rows:
                sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        return "\n\n".join(sheets) if sheets else f"[XLSX: {filename} — no data found]"
    except ImportError:
        logger.warning("openpyxl not installed; cannot extract .xlsx text")
        return f"[XLSX: {filename} — install openpyxl to enable text extraction]"
    except Exception as exc:
        logger.warning("XLSX extraction error for %s: %s", filename, exc)
        return f"[XLSX: {filename} — extraction error: {exc}]"


def _extract_pptx(data: bytes, filename: str) -> str:
    try:
        from pptx import Presentation  # noqa: F401
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else f"[PPTX: {filename} — no text found]"
    except ImportError:
        logger.warning("python-pptx not installed; cannot extract .pptx text")
        return f"[PPTX: {filename} — install python-pptx to enable text extraction]"
    except Exception as exc:
        logger.warning("PPTX extraction error for %s: %s", filename, exc)
        return f"[PPTX: {filename} — extraction error: {exc}]"


def _extract_zip(data: bytes, filename: str) -> str:
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        return f"[ZIP: {filename} — invalid or corrupted archive]"

    parts = [f"[ZIP archive: {filename}]", "Contents:"]
    members = zf.namelist()
    parts.append(", ".join(members[:50]) + ("..." if len(members) > 50 else ""))
    parts.append("")

    # R9.4 (X8) AC2: guard against a zip bomb using the archive DIRECTORY's declared
    # uncompressed sizes (ZipInfo.file_size) — checked BEFORE any member is read, so
    # a 10 KB archive that expands to gigabytes is never decompressed. Both a
    # per-member cap and a running total cap apply.
    info_by_name = {i.filename: i for i in zf.infolist()}
    text_members = [m for m in members if Path(m).suffix.lower() in TEXT_EXTENSIONS]
    total_uncompressed = 0
    for member in text_members[:10]:  # cap at 10 text files
        info = info_by_name.get(member)
        declared = info.file_size if info else 0
        if declared > MAX_ZIP_MEMBER_BYTES:
            parts.append(f"=== {member} === [too large to display]")
            continue
        total_uncompressed += declared
        if total_uncompressed > MAX_ZIP_TOTAL_BYTES:
            parts.append(f"=== {member} === [archive decompression limit reached]")
            break
        try:
            member_data = zf.read(member)
            if len(member_data) > 200_000:
                text = f"[{member} — too large to display]"
            else:
                text = _read_plain_text(member_data, member)[:5000]
            parts.append(f"=== {member} ===\n{text}")
        except Exception:
            parts.append(f"=== {member} === [unreadable]")

    return "\n".join(parts)


def _extract_text(data: bytes, filename: str, suffix: str) -> str:
    suffix = suffix.lower()

    if suffix in TEXT_EXTENSIONS:
        return _read_plain_text(data, filename)

    if suffix == ".pdf":
        return _extract_pdf(data, filename)

    if suffix == ".docx":
        return _extract_docx(data, filename)

    if suffix in (".xlsx", ".xls"):
        return _extract_xlsx(data, filename)

    if suffix == ".pptx":
        return _extract_pptx(data, filename)

    if suffix == ".doc":
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs) if paragraphs else f"[DOC: {filename} — no text found]"
        except Exception:
            return f"[DOC: {filename} — old .doc format could not be read. Please save as .docx and re-upload.]"

    if suffix == ".zip":
        return _extract_zip(data, filename)

    if suffix in IMAGE_EXTENSIONS:
        import base64
        mime_map = {
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".webp": "image/webp",
            ".heic": "image/heic", ".bmp": "image/bmp", ".tiff": "image/tiff",
        }
        mime = mime_map.get(suffix, "image/jpeg")
        b64 = base64.b64encode(data).decode()
        return f"__IMAGE_DATA__data:{mime};base64,{b64}__FILENAME__{filename}"

    if suffix in MEDIA_EXTENSIONS:
        return (
            f"[Audio/Video: {filename}]\n"
            "Note: Media files cannot be transcribed in this version."
        )

    return f"[File: {filename} — unsupported format '{suffix}']"


def may_read_images(user: dict) -> bool:
    """Who may have Flo read a photograph.

    Settled by Abhimanyu, 2026-07-22 (revised the same day): **Owner, Principal and
    the other office staff — accountant, receptionist and the rest of the admin
    roles. NOT teachers, and NOT students.**

    The revision is the point. An earlier draft included teachers on the reasoning
    that they photograph forms; the owner narrowed it to the office. This is
    paperwork handling — fee slips, admission forms, circulars — and it belongs with
    the people who do the paperwork. Teachers were removed deliberately, so do not
    "restore" them on the assumption it was an oversight.

    Deliberately NOT `is_owner_or_principal` either, which would exclude the
    accountant and receptionist who handle most of the paper.
    """
    role = (user or {}).get("role")
    return role in ("owner", "admin")


@router.post("/upload")
async def upload_chat_file(request: Request, file: UploadFile = File(...)):
    """Accept a file and return extracted text for inclusion in the AI conversation context."""
    user = get_current_user(request)  # auth guard — raises 401 if not authenticated

    filename = file.filename or "uploaded_file"
    suffix = Path(filename).suffix.lower()
    if suffix in BLOCKED_EXTENSIONS:
        raise HTTPException(415, f"File type {suffix} is not permitted")

    # R9.4 (X8) AC1: reject early on the declared Content-Length, then read in
    # bounded chunks and abort the moment the cap is exceeded — so an oversized
    # (or lying-about-its-size) upload is never fully buffered into memory.
    content_length = request.headers.get("content-length")
    if content_length and content_length.isdigit() and int(content_length) > MAX_FILE_SIZE_BYTES:
        raise HTTPException(413, f"File too large (max {MAX_FILE_SIZE_BYTES // (1024*1024)} MB)")

    buf = bytearray()
    while True:
        chunk = await file.read(_READ_CHUNK_BYTES)
        if not chunk:
            break
        buf.extend(chunk)
        if len(buf) > MAX_FILE_SIZE_BYTES:
            raise HTTPException(413, f"File too large (max {MAX_FILE_SIZE_BYTES // (1024*1024)} MB)")
    data = bytes(buf)

    extracted = _extract_text(data, filename, suffix)

    # Images return a special sentinel with the base64 data URL
    image_data = None
    ocr_note = None
    used_paid_vision = False
    if extracted.startswith("__IMAGE_DATA__"):
        parts = extracted.split("__FILENAME__")
        image_data = parts[0].replace("__IMAGE_DATA__", "")

        # Epic 10 / Story 10.5: read the words off the page HERE, on this server.
        # Doing it at the upload boundary means Flo receives ordinary text and the
        # whole chat pipeline needs no knowledge of images at all. Free, and the
        # picture never leaves this machine.
        if not may_read_images(user):
            extracted = (
                f"[Image attached: {filename}. Reading images is available to the "
                "Owner, the Principal and office staff only, so its contents were not read.]"
            )
            image_data = None  # do not hand the bytes on to a caller who may not use them
        else:
            result = extract_image_text(data)
            if not result.available:
                # NOT an empty page. Say which, or a deployment problem reads as
                # "this form is blank" — the Epic 4 defect in a new place.
                ocr_note = result.reason
                extracted = f"[Image attached: {filename}. {result.reason}]"
            elif result.found_text:
                extracted = (
                    f"[Text read from the image {filename}"
                    + (f", language {result.language}" if result.language else "")
                    + "]\n" + result.text
                )
                if result.notes:
                    extracted += "\n\n[" + " ".join(result.notes) + "]"
            else:
                # Story 10.6: reading the words was not enough, so fall back to the
                # paid service — and ONLY here. A page whose text was read never
                # reaches this branch, which is what keeps printed paper free.
                vision = describe_image_bytes(data, sniff_image_type(data) or "image/jpeg")
                if vision.understood:
                    used_paid_vision = True
                    extracted = f"[Description of the image {filename}]\n{vision.description}"
                else:
                    ocr_note = vision.reason or result.reason
                    extracted = f"[Image attached: {filename}. {ocr_note}]"

            # Reading an image is a read of school records and may involve a child.
            # Ids and counts only (NFR-S2) — never the text that was read.
            try:
                await write_audit(
                    get_db(),
                    action="image_text_read",
                    entity_id=filename[:120],
                    collection="chat_uploads",
                    changed_by=user.get("id", ""),
                    changed_by_role=user.get("role", ""),
                    school_id=get_school_id(),
                    branch_id=user.get("branch_id", ""),
                    changes={
                        "size_bytes": len(data),
                        "chars_found": result.char_count,
                        "engine_available": result.available,
                        # Recorded so the owner can see how often the PAID path was
                        # taken, rather than discovering it on a bill.
                        "paid_vision_used": used_paid_vision,
                    },
                )
            except Exception:
                logger.warning("image_text_read audit failed", exc_info=True)

    if not image_data and len(extracted) > MAX_TEXT_LENGTH:
        extracted = extracted[:MAX_TEXT_LENGTH] + f"\n\n[... truncated — original {len(extracted):,} chars]"

    return {
        "success": True,
        "filename": filename,
        "size_bytes": len(data),
        "extracted_text": extracted,
        "char_count": len(extracted),
        "image_data": image_data,
        # Present only when an image could not be read, so the screen can say why
        # instead of showing an attachment that silently contributed nothing.
        "ocr_note": ocr_note,
    }
