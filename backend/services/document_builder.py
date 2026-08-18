"""Turn a structured description into a real Word, Excel, PowerPoint or PDF file.

UI Sweep Epic 10, Story 10.1.

WHY THIS EXISTS. Flo told the owner it could write the *content* of a circular or a fee
sheet but "not directly generate a real .docx file in this setup". That was Flo
underselling the platform: `python-docx`, `openpyxl`, `python-pptx` and `fpdf2` are all
already pinned in requirements.txt. Three of them were only ever used to READ uploaded
files (`routes/chat_upload.py`); only PDF was used to write, for certificates and fee
receipts. Nothing was missing except a place to put the writing code.

THE RULE THIS MODULE ENFORCES: there is ONE builder. Four half-built generators
scattered across route files is what happens otherwise, and then a fix to page size or
filename sanitising lands in one of them.

WHAT THIS MODULE IS NOT: it does not fetch data, decide who may see it, store anything,
or write an audit row. It takes a description and returns bytes. Authorization belongs
to the caller, because the caller knows which data it drew on - see
`services/document_export.py` for the storing half and `ai/tool_functions_v2.py` for the
gate.
"""

from __future__ import annotations

import io
import logging
import os
import re
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# Caps mirror the `.to_list(N)` limits already used in routes/exports.py. A document
# nobody can open is not a successful export, and an unbounded row count is how a
# request for "every fee record" becomes an out-of-memory on a small instance.
MAX_ROWS = 5000
MAX_COLUMNS = 60
MAX_CELL_CHARS = 2000
MAX_SLIDES = 100
MAX_PARAGRAPHS = 2000

CONTENT_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "csv": "text/csv",
    "xml": "application/xml",
    "md": "text/markdown",
    "txt": "text/plain",
}

SUPPORTED_TYPES = tuple(CONTENT_TYPES)

# Formats that NEVER carry the school's branding, whatever the caller asks for.
# A spreadsheet is data: rows of school name above the column headings shift every
# row down, so a formula, a filter or an import into another system starts on the
# wrong line. Every other format is something a person reads, where the branding is
# the point.
# XML is also unbranded: it is a data-interchange format and adding school header
# rows would break any parser reading it.
UNBRANDED_TYPES = frozenset({"xlsx", "csv", "xml"})


class DocumentBuildError(Exception):
    """The description could not be turned into a file.

    Raised BEFORE anything is stored, so a malformed request never leaves an orphan
    object in S3 with no `file_uploads` record pointing at it.
    """


@dataclass
class BuiltDocument:
    content: bytes
    content_type: str
    filename: str
    doc_type: str
    truncated: bool = False
    notes: List[str] = field(default_factory=list)
    # The same document as simple HTML, so a person can read it, correct a sentence
    # and download the corrected copy. Built here from the description rather than by
    # parsing the finished .pdf or .docx back apart, which is unreliable and would be
    # a second definition of what the document says. See `editable_html()`.
    editable_html: str = ""

    @property
    def size_bytes(self) -> int:
        return len(self.content)


# ── Filename hygiene ────────────────────────────────────────────────────────────

_UNSAFE = re.compile(r"[^A-Za-z0-9._-]+")


def safe_filename(name: str, doc_type: str) -> str:
    """Make a filename safe for an S3 key AND for a Content-Disposition header.

    Both matter. A path separator walks outside the intended `{school_id}/uploads/`
    prefix; a newline or a quote forges a response header. The name may come from a
    person or from Flo, so neither is trusted.
    """
    base = (name or "").strip()
    # Take the last path component before sanitising, so "../../etc/passwd" cannot
    # survive as "......etc-passwd" and read like a traversal attempt.
    base = base.replace("\\", "/").split("/")[-1]
    for suffix in SUPPORTED_TYPES:
        if base.lower().endswith("." + suffix):
            base = base[: -(len(suffix) + 1)]
            break
    base = _UNSAFE.sub("-", base).strip("-._")
    base = re.sub(r"-{2,}", "-", base)
    if not base:
        base = "document"
    return f"{base[:80]}.{doc_type}"


def _clean_cell(value: Any) -> str:
    if value is None:
        return ""
    text = str(value)
    return text[:MAX_CELL_CHARS]


def _validate(doc_type: str, title: str) -> None:
    if doc_type not in SUPPORTED_TYPES:
        raise DocumentBuildError(
            f"Unsupported document type '{doc_type}'. "
            f"Supported: {', '.join(SUPPORTED_TYPES)}."
        )
    if title is not None and len(str(title)) > 300:
        raise DocumentBuildError("Title is too long (max 300 characters).")


def _normalise_table(headers: Optional[List[Any]], rows: Optional[List[List[Any]]],
                     max_rows: Optional[int] = None):
    """Return (headers, rows, truncated). Ragged rows are padded, not rejected -
    real data is ragged, and refusing the whole document over one short row would
    be worse than filling a blank.

    `max_rows` is read at CALL time, not bound as a default. Binding it would freeze
    MAX_ROWS at import and quietly break every test and every caller that changes the
    constant - which is exactly what happened when this argument was first added."""
    max_rows = MAX_ROWS if max_rows is None else max_rows
    hdrs = [_clean_cell(h) for h in (headers or [])][:MAX_COLUMNS]
    raw = rows or []
    truncated = len(raw) > max_rows
    out: List[List[str]] = []
    width = len(hdrs)
    for row in raw[:max_rows]:
        if not isinstance(row, (list, tuple)):
            row = [row]
        cells = [_clean_cell(c) for c in row][:MAX_COLUMNS]
        width = max(width, len(cells))
        out.append(cells)
    if not hdrs and out:
        width = max(len(r) for r in out)
    for row in out:
        row.extend([""] * (width - len(row)))
    if hdrs:
        hdrs.extend([""] * (width - len(hdrs)))
    return hdrs, out, truncated


# ── Hindi in PDFs ───────────────────────────────────────────────────────────────
#
# THE DEFECT THIS CLOSES. fpdf2's built-in fonts are Latin-1 only. Every Devanagari
# character was therefore turned into "?" by `_latin1()` below, silently. Word, Excel,
# PowerPoint, CSV, Markdown and plain text all handled Hindi correctly the whole time,
# so this was PDF alone. Losing the text quietly was the wrong trade for a school in
# Amroha: a circular to parents that reads "?????? ??????" is not a document.
#
# THE FIX. Register Noto Sans Devanagari, a real Unicode font, and use it instead of
# Helvetica. It also covers Latin, so English documents are unaffected in content and
# only change in typeface.
#
# LICENCE. SIL Open Font License 1.1, checked by reading the licence text shipped
# beside the font rather than taking it on trust. It grants permission to "use, study,
# copy, merge, embed, modify, redistribute" - embedding in a generated PDF and
# shipping the file with the server are both covered. Its one relevant condition is
# that the licence travels with the font, so `assets/fonts/OFL.txt` sits next to the
# .ttf files and `.ebignore` carries an explicit exception to keep it in the deploy.
#
# SHAPING IS NOT OPTIONAL FOR HINDI, and this is the part that is easy to get wrong.
# Devanagari is not written in the order it is stored: the vowel sign in "हि" is typed
# after the consonant but drawn before it, and consonant clusters join into single
# conjunct glyphs. Dropping the codepoints onto the page in storage order produces
# text made of the right letters in visibly wrong shapes and order. So text shaping is
# turned on when `uharfbuzz` is installed, and that is what makes the output correct
# rather than merely non-question-marks.
#
# THREE TIERS, DEGRADING NEVER FAILING. A missing font file or a missing shaper must
# not turn every PDF in the school into an error - that would be a far bigger outage
# than the defect being fixed:
#   1. font + shaping  -> correct Hindi. This is the shipped configuration.
#   2. font, no shaper -> Hindi characters appear but conjuncts and vowel signs are
#                         misplaced. Wrong-looking, but the words are recoverable.
#   3. no font         -> exactly today's behaviour, Latin-1 with "?" substitution.
# Which tier was used is written to the log, because a silent drop to tier 3 on the
# live server is how this defect stayed invisible in the first place.

FONT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets", "fonts")

# Resolved from THIS file's location, never the working directory. Elastic Beanstalk
# starts the app from a different cwd than local development does, and a relative path
# here is how the font would be found in testing and missing in production.
# TWO families, not one, and this was found by a test rather than by reasoning.
#
# The first attempt used Noto Sans Devanagari alone, on the assumption that it covered
# Latin as well. It does not. English text came out of the PDF completely BLANK: the
# subsetter kept only "space" and "period" out of "Holiday on Monday." - every letter
# was dropped. That would have replaced a small Hindi defect with a total English one,
# on a platform where nearly every document is in English.
#
# So Noto Sans (Latin) is the document font, and Noto Sans Devanagari is registered as
# its FALLBACK. fpdf2 reaches for the fallback per character, when the main font has no
# glyph, which is exactly the behaviour wanted for a mixed sentence like
# "कक्षा III की उपस्थिति".
UNICODE_FONT_FAMILY = "NotoSans"
DEVANAGARI_FONT_FAMILY = "NotoSansDevanagari"

UNICODE_FONT_FILES = {
    "": os.path.join(FONT_DIR, "NotoSans-Regular.ttf"),
    "B": os.path.join(FONT_DIR, "NotoSans-Bold.ttf"),
}
DEVANAGARI_FONT_FILES = {
    "": os.path.join(FONT_DIR, "NotoSansDevanagari-Regular.ttf"),
    "B": os.path.join(FONT_DIR, "NotoSansDevanagari-Bold.ttf"),
}

# Any character above Latin-1 means the core fonts cannot render the text.
_NON_LATIN1 = re.compile(r"[^\x00-\xff]")


def unicode_font_available() -> bool:
    """Are BOTH font families actually on disk? Checked, never assumed.

    Both, because the Latin font alone would silently drop Hindi and the Devanagari
    font alone would silently drop English. Either half missing means fall back.
    """
    paths = list(UNICODE_FONT_FILES.values()) + list(DEVANAGARI_FONT_FILES.values())
    return all(os.path.exists(path) for path in paths)


def text_shaping_available() -> bool:
    """Is the shaper installed? Without it Hindi renders in the wrong shapes."""
    try:
        import uharfbuzz  # noqa: F401
    except ImportError:
        return False
    return True


def _register_unicode_font(pdf) -> bool:
    """Attach the Devanagari font to this document. False means fall back.

    Deliberately swallows every failure. A corrupt or unreadable font file must
    degrade to the old Latin-1 behaviour, not take down document generation for the
    whole school.
    """
    if not unicode_font_available():
        logger.warning(
            "PDF: Devanagari font not found at %s - falling back to Latin-1, so any "
            "Hindi in this document will be replaced with '?'", FONT_DIR,
        )
        return False
    try:
        for style, path in UNICODE_FONT_FILES.items():
            pdf.add_font(UNICODE_FONT_FAMILY, style, path)
        for style, path in DEVANAGARI_FONT_FILES.items():
            pdf.add_font(DEVANAGARI_FONT_FAMILY, style, path)
        # Per-character fallback: the Latin font draws the English, and anything it has
        # no glyph for is handed to the Devanagari font. This is what makes a mixed
        # sentence work rather than forcing a whole document into one script.
        pdf.set_fallback_fonts([DEVANAGARI_FONT_FAMILY])
    except Exception:  # pragma: no cover - defensive, see docstring
        logger.exception("PDF: could not register the Unicode fonts - falling back to Latin-1")
        return False

    if text_shaping_available():
        try:
            pdf.set_text_shaping(True)
        except Exception:  # pragma: no cover
            logger.warning("PDF: text shaping could not be enabled; Hindi will render unshaped")
    else:
        logger.warning(
            "PDF: uharfbuzz is not installed, so Devanagari will render unshaped - "
            "the characters are correct but conjuncts and vowel signs sit wrong"
        )
    return True


def _build_pdf_text_encoder(unicode_ok: bool):
    """Return the function that prepares a string for the page.

    With the Unicode font, text passes through untouched. Without it, the original
    Latin-1 substitution is used, which is lossy and known to be so.
    """
    if unicode_ok:
        return lambda text: str(text)

    def _latin1(text: str) -> str:
        # The pre-2026-08-07 behaviour, kept ONLY as the last-resort tier. Devanagari
        # would otherwise raise and lose the whole document, so it is replaced rather
        # than crashing.
        return str(text).encode("latin-1", "replace").decode("latin-1")

    return _latin1


# ── The same document, as editable HTML ─────────────────────────────────────────
#
# WHY. Until now a document Flo made could only be downloaded. There was no way to
# read one, fix a sentence and take the corrected copy - you had to ask Flo again and
# hope. Abhimanyu asked for an edit panel on 2026-08-07 and decided the corrected copy
# is DOWNLOADED ONLY and nothing is saved back to the server, so this is purely a
# starting point for the browser editor. Nothing here is ever stored as the document.
#
# Built from the description, NOT by reading the finished PDF back apart. Parsing a
# PDF into editable text is unreliable, and it would make the file a second, competing
# statement of what the document says.

def _escape(text: Any) -> str:
    """HTML-escape. The content comes from Flo and from school data, so it is never
    trusted as markup - the browser side sanitises again on top of this."""
    return (
        str(text)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def editable_html(title, paragraphs, headers, rows, truncated_note) -> str:
    """The document as plain HTML a contentEditable panel can hold."""
    parts: List[str] = []
    if title:
        parts.append(f"<h1>{_escape(str(title)[:300])}</h1>")
    for para in (paragraphs or [])[:MAX_PARAGRAPHS]:
        parts.append(f"<p>{_escape(_clean_cell(para))}</p>")
    if rows:
        parts.append('<table border="1" cellspacing="0" cellpadding="6">')
        if headers:
            head = "".join(f"<th>{_escape(h)}</th>" for h in headers)
            parts.append(f"<thead><tr>{head}</tr></thead>")
        body = "".join(
            "<tr>" + "".join(f"<td>{_escape(cell)}</td>" for cell in row) + "</tr>"
            for row in rows
        )
        parts.append(f"<tbody>{body}</tbody></table>")
    if truncated_note:
        parts.append(f"<p><em>{_escape(truncated_note)}</em></p>")
    return "".join(parts)


# ── The school's letterhead ─────────────────────────────────────────────────────
#
# WHY. Flo's documents were a bare title, bare paragraphs and an unstyled grid. The
# school sends these to parents and to the board, and a page with no letterhead does
# not read as a school document at all. Abhimanyu asked for this on 2026-08-07 after
# comparing Flo's output with a hand-made one.
#
# WHERE THE DESIGN COMES FROM. It is copied from the school's OWN printed enquiry
# form (`aaryans_database/`): logo on the left, the school name across the top, the
# CBSE affiliation line beneath it, a rule, and the address and contact details in a
# footer under a second rule. Nothing here is invented, which matters because a made
# up affiliation number on a document sent to parents is worse than no letterhead.
#
# The text comes from `school_identity.py`, the one verified source, so a correction
# there reaches these documents without a second edit. It reads the module constants
# only - no database call - so this module keeps its promise of taking a description
# and returning bytes.

_ASSET_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets")

# The crest at the top of the letterhead: full colour, transparent background.
LOGO_PATH = os.path.join(_ASSET_DIR, "aaryans-logo.png")

# The SAME crest, desaturated and lightened, for the background watermark. A separate
# file rather than a live filter: full colour behind a fee table reads as a mistake,
# while the printed form shows a pale, nearly colourless ghost. Falls back to the
# ordinary logo if it is missing, so the watermark never simply disappears.
CREST_WATERMARK_PATH = os.path.join(_ASSET_DIR, "aaryans-crest-watermark.png")

# From the school's printed form. Navy, matching the header type and the rules on it.
BRAND_NAVY = (22, 43, 90)
BRAND_BLUE = (31, 106, 184)
RULE_GREY = (210, 218, 232)

# The tiled wordmark, tuned against a zoomed photograph of the printed form rather
# than guessed. Kept as named constants because this is the part most likely to need
# nudging by eye: step sizes are millimetres on A4 (210mm wide), so 26mm gives about
# eight repetitions across the page, which is what the paper shows.
WATERMARK_STEP_X = 26.0
# 14 rows down the sheet, close together with a small gap, counted off the printed
# form. Expressed as a row COUNT rather than a spacing so it stays right if the page
# size ever changes; the spacing is derived from the page height.
WATERMARK_ROWS = 14
WATERMARK_FONT_PT = 7
WATERMARK_INK = (120, 145, 185)
WATERMARK_OPACITY = 0.13
# The chariot sits behind everything, larger and fainter than the wordmark.
CHARIOT_OPACITY = 0.07
CHARIOT_SIZE_MM = 120


def _identity() -> Dict[str, str]:
    """The school's verified details. Falls back to empty strings, never invents."""
    try:
        from school_identity import default_school_identity

        return default_school_identity()
    except Exception:  # pragma: no cover - identity module is always present
        logger.exception("PDF: could not read the school identity; letterhead will be plain")
        return {}


def _stationery() -> Dict[str, str]:
    """How the school's details are PRINTED on its own paper.

    Separate from `_identity()` on purpose. The identity values are what the product
    shows on screen; these are the exact wording and punctuation of the printed
    enquiry form, which is what a document pretending to be school stationery has to
    reproduce. Both live in `school_identity.py`, so there is still one file to edit.
    """
    try:
        from school_identity import LETTERHEAD

        return dict(LETTERHEAD)
    except Exception:  # pragma: no cover
        logger.exception("PDF: could not read the printed letterhead wording")
        return {}


def letterhead_lines() -> List[str]:
    """The letterhead as plain lines, for the formats that have no page furniture.

    Excel, PowerPoint, Markdown, plain text and CSV cannot carry a header and footer
    the way a PDF or a Word file can, but Abhimanyu's instruction on 2026-08-07 was
    that this is the school's branding on EVERY type of document, not only the two
    that happen to have page margins. So those formats get the same wording written
    into the content itself.
    """
    s = _stationery()
    return [line for line in (
        s.get("name", ""),
        s.get("tagline", ""),
        s.get("affiliation_line", ""),
    ) if line]


def letterhead_footer_lines() -> List[str]:
    """The address block, for the same formats."""
    s = _stationery()
    return [line for line in (
        s.get("footer_address", ""),
        s.get("footer_contact", ""),
    ) if line]


def _letterhead_pdf_class(base_cls, identity: Dict[str, str]):
    """Build an FPDF subclass that draws the letterhead on EVERY page.

    A subclass rather than drawing it once, because `header()` and `footer()` are
    called by fpdf2 on each new page. Drawing it inline would put the letterhead on
    page one of a twelve-page fee report and nothing on the rest, which is exactly how
    a document stops looking official halfway through.

    The font family and the text encoder are read off the INSTANCE (`self._family`,
    `self._enc`) rather than captured here, because which font is in use is only known
    after registration is attempted, and registration needs the instance to exist
    first. Capturing them here would have baked in the guess made before the fonts
    were actually loaded.
    """
    # Every string below is the PRINTED wording, taken from the school's own enquiry
    # form. It is deliberately not rebuilt out of the identity fields: doing that is
    # what produced "The Aaryans" where the paper says "THE AARYANS", and a reworded
    # footer where the paper has an exact address line (Abhimanyu, 2026-08-07).
    stationery = _stationery()
    school = stationery.get("name") or identity.get("school_name", "")
    tagline = stationery.get("tagline", "")
    affiliation_line = stationery.get("affiliation_line", "")
    footer_address = stationery.get("footer_address", "")
    footer_contact = stationery.get("footer_contact", "")
    watermark_text = stationery.get("watermark_text", "")

    class Letterhead(base_cls):
        # Safe defaults so a page drawn before the fonts are settled still renders.
        _family = "Helvetica"
        _enc = staticmethod(lambda text: str(text).encode("latin-1", "replace").decode("latin-1"))

        def _watermark(self):
            """The pale repeated wordmark and chariot printed across the whole page.

            Called from `footer()`, so it overlays the finished page. See the note
            there for why it cannot sit underneath.

            Kept very faint on purpose. A watermark that competes with the text turns
            a readable fee sheet into a difficult one, and these are read on paper and
            on phones.
            """
            family, enc = self._family, self._enc
            # Automatic page breaks are turned OFF for the duration. The tiling loop
            # writes near the bottom of the sheet, which tripped the page break, which
            # called header() again, which drew another watermark: the first version
            # of this recursed until Python gave up. The watermark is decoration and
            # must never create a page.
            auto = self.auto_page_break
            break_margin = self.b_margin
            self.set_auto_page_break(False)
            try:
                # Two passes, because the two elements are not equally strong on the
                # paper: the chariot is a ghost behind everything, the wordmark is
                # faint but clearly legible.
                crest_path = (
                    CREST_WATERMARK_PATH if os.path.exists(CREST_WATERMARK_PATH) else LOGO_PATH
                )
                if os.path.exists(crest_path):
                    with self.local_context(fill_opacity=CHARIOT_OPACITY, stroke_opacity=CHARIOT_OPACITY):
                        size = CHARIOT_SIZE_MM
                        self.image(
                            crest_path,
                            x=(self.w - size) / 2,
                            y=(self.h - size) / 2,
                            w=size,
                        )

                with self.local_context(fill_opacity=WATERMARK_OPACITY, stroke_opacity=WATERMARK_OPACITY):
                    # The tiled wordmark. Rows are offset by half a step so the grid
                    # reads as a texture rather than as columns.
                    if watermark_text:
                        # Matched against a zoomed photograph of the printed form
                        # (Abhimanyu, 2026-08-07): WATERMARK_ROWS rows down the sheet
                        # with only a small gap between them, in a light blue.
                        self.set_font(family, "B", WATERMARK_FONT_PT)
                        self.set_text_color(*WATERMARK_INK)
                        step_x = WATERMARK_STEP_X
                        step_y = self.h / WATERMARK_ROWS
                        line_h = 4.0

                        # The chariot's box, so the wording can BREAK AROUND IT rather
                        # than run across it. On the printed form the crest shows
                        # through cleanly and the repeated text stops at its edge and
                        # picks up again on the other side; printing over it turned the
                        # crest into a smudge (Abhimanyu, with a screenshot).
                        chariot = None
                        if os.path.exists(LOGO_PATH):
                            size = CHARIOT_SIZE_MM
                            cx, cy = (self.w - size) / 2, (self.h - size) / 2
                            # A small margin so the text does not touch the crest edge.
                            pad = 3.0
                            chariot = (cx - pad, cy - pad, cx + size + pad, cy + size + pad)

                        def _clear_of_chariot(x0, y0):
                            """True when this tile does not overlap the crest at all."""
                            if chariot is None:
                                return True
                            left, top, right, bottom = chariot
                            x1, y1 = x0 + step_x, y0 + line_h
                            return x1 <= left or x0 >= right or y1 <= top or y0 >= bottom

                        row = 0
                        y = 2.0
                        while y + line_h <= self.h:
                            # Never start at a negative x: fpdf2 refuses to draw a
                            # cell that begins off the page, and the whole watermark
                            # was being lost to that one detail.
                            x = (step_x / 2) if row % 2 else 0.0
                            while x + step_x <= self.w:
                                if _clear_of_chariot(x, y):
                                    self.set_xy(x, y)
                                    self.cell(step_x, line_h, enc(watermark_text), align="C")
                                x += step_x
                            y += step_y
                            row += 1
            except Exception:  # pragma: no cover - decoration must never cost the document
                logger.exception("PDF: the watermark could not be drawn")
            finally:
                self.set_auto_page_break(auto, margin=break_margin)
            self.set_text_color(0, 0, 0)

        def header(self):
            family, enc = self._family, self._enc
            top = 8
            if os.path.exists(LOGO_PATH):
                try:
                    self.image(LOGO_PATH, x=self.l_margin, y=top, h=18)
                except Exception:  # pragma: no cover - a bad image must not stop the doc
                    logger.warning("PDF: the school logo could not be drawn")

            # Name and tagline are centred across the full width, like the printed
            # form, rather than being pushed right by the logo.
            self.set_xy(self.l_margin, top)
            self.set_font(family, "B", 26)
            self.set_text_color(*BRAND_NAVY)
            self.cell(0, 11, enc(school), align="C")

            self.set_xy(self.l_margin, top + 11)
            self.set_font(family, size=8)
            self.set_text_color(60, 72, 96)
            self.cell(0, 4, enc(tagline), align="C")
            if affiliation_line:
                self.set_xy(self.l_margin, top + 15)
                self.cell(0, 4, enc(affiliation_line), align="C")

            rule_y = top + 21
            self.set_draw_color(*BRAND_BLUE)
            self.set_line_width(0.6)
            self.line(self.l_margin, rule_y, self.w - self.r_margin, rule_y)

            self.set_text_color(0, 0, 0)
            self.set_line_width(0.2)
            self.set_y(rule_y + 6)

        def footer(self):
            family, enc = self._family, self._enc

            # The watermark is drawn HERE, not in header(), and that is the whole
            # point. fpdf2 calls footer() after the page's body has been written, so
            # anything drawn here lands ON TOP.
            #
            # Drawing it in header() put it underneath, and the table then painted
            # over it: banded rows and the heading row use opaque fills, so the
            # chariot vanished behind every table (Abhimanyu, 2026-08-07, with a
            # screenshot). Nothing could fix that from underneath short of giving up
            # the row banding.
            #
            # Overlaying is also closer to the printed form, where the wordmark is on
            # the paper and shows across the whole sheet. It is only safe because the
            # opacities are low enough that the text underneath stays fully readable.
            self._watermark()

            self.set_y(-16)
            self.set_draw_color(*BRAND_BLUE)
            self.set_line_width(0.5)
            self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())

            # Both lines exactly as printed on the school's form, punctuation and all.
            self.set_y(-13)
            self.set_font(family, size=7)
            self.set_text_color(60, 72, 96)
            self.cell(0, 3.5, enc(footer_address), align="C")

            self.set_y(-9.5)
            self.cell(0, 3.5, enc(footer_contact), align="C")

            # Page numbers matter on anything the school files or posts: a loose sheet
            # from a twelve-page fee report is unidentifiable without them.
            self.set_y(-6)
            self.set_font(family, size=7)
            self.cell(0, 3.5, enc(f"Page {self.page_no()} of {{nb}}"), align="C")
            self.set_text_color(0, 0, 0)

    return Letterhead


# ── Builders, one per type ──────────────────────────────────────────────────────

def _docx_letterhead(doc) -> None:
    """Put the school's header and footer into a Word document's page furniture.

    Written into the SECTION header and footer, not as ordinary paragraphs at the top
    of the body, so Word repeats them on every page and a person editing the circular
    cannot delete half the letterhead by accident while typing.
    """
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    # Same printed wording as the PDF letterhead, from the one place that holds it.
    stationery = _stationery()
    section = doc.sections[0]

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if os.path.exists(LOGO_PATH):
        try:
            header.add_run().add_picture(LOGO_PATH, height=Pt(34))
            header.add_run("  ")
        except Exception:  # pragma: no cover - a bad image must not stop the doc
            logger.warning("DOCX: the school logo could not be placed")

    name_run = header.add_run(stationery.get("name", ""))
    name_run.bold = True
    name_run.font.size = Pt(22)
    name_run.font.color.rgb = RGBColor(*BRAND_NAVY)

    tagline = header.add_run(
        f"\n{stationery.get('tagline', '')}"
        f"\n{stationery.get('affiliation_line', '')}"
    )
    tagline.font.size = Pt(8)
    tagline.font.color.rgb = RGBColor(60, 72, 96)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    contact = footer.add_run(
        f"{stationery.get('footer_address', '')}\n{stationery.get('footer_contact', '')}"
    )
    contact.font.size = Pt(7)
    contact.font.color.rgb = RGBColor(60, 72, 96)


def _build_docx(title, paragraphs, headers, rows, truncated_note, letterhead=True):
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - dependency is pinned
        raise DocumentBuildError("Word support is not installed on this server.") from exc

    doc = Document()
    if letterhead:
        try:
            _docx_letterhead(doc)
        except Exception:  # pragma: no cover - never lose the document over decoration
            logger.exception("DOCX: the letterhead could not be built; carrying on without it")
    if title:
        doc.add_heading(str(title)[:300], level=1)
    for para in (paragraphs or [])[:MAX_PARAGRAPHS]:
        doc.add_paragraph(_clean_cell(para))
    if rows:
        table = doc.add_table(rows=1 if headers else 0, cols=len(rows[0]) or 1)
        table.style = "Table Grid"
        if headers:
            for i, h in enumerate(headers):
                table.rows[0].cells[i].text = h
        for row in rows:
            cells = table.add_row().cells
            for i, value in enumerate(row):
                cells[i].text = value
    if truncated_note:
        doc.add_paragraph(truncated_note)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx(title, paragraphs, headers, rows, truncated_note, letterhead=True):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except ImportError as exc:  # pragma: no cover
        raise DocumentBuildError("Excel support is not installed on this server.") from exc

    wb = Workbook()
    ws = wb.active
    # Sheet names cannot exceed 31 chars or contain []:*?/\ - Excel refuses to open
    # the file rather than complaining, so this is silently fatal if not handled.
    ws.title = (re.sub(r"[\[\]:*?/\\]", "-", str(title or "Sheet"))[:31]) or "Sheet"

    if letterhead:
        # Written into the sheet AND into the print header, so it is the school's
        # document both on screen and on paper.
        for index, line in enumerate(letterhead_lines()):
            ws.append([line])
            cell = ws.cell(row=ws.max_row, column=1)
            cell.font = Font(bold=(index == 0), size=14 if index == 0 else 9,
                             color="FF162B5A" if index == 0 else "FF3C4860")
        ws.append([])
        ws.oddHeader.center.text = "&B" + (letterhead_lines() or [""])[0]
        ws.oddFooter.center.text = "\n".join(letterhead_footer_lines())

    for para in (paragraphs or [])[:20]:
        ws.append([_clean_cell(para)])
    if paragraphs:
        ws.append([])
    if headers:
        ws.append(headers)
        for cell in ws[ws.max_row]:
            cell.font = Font(bold=True)
    for row in rows:
        ws.append(row)
    if truncated_note:
        ws.append([])
        ws.append([truncated_note])

    if letterhead:
        ws.append([])
        for line in letterhead_footer_lines():
            ws.append([line])
            ws.cell(row=ws.max_row, column=1).font = Font(size=8, color="FF3C4860")

    # Width by content so a fee sheet opens readable rather than as ####.
    for idx, column_cells in enumerate(ws.columns, start=1):
        longest = max((len(str(c.value)) for c in column_cells if c.value is not None), default=0)
        ws.column_dimensions[ws.cell(row=1, column=idx).column_letter].width = min(60, max(10, longest + 2))

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx(title, paragraphs, headers, rows, truncated_note, slides, letterhead=True):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover
        raise DocumentBuildError("PowerPoint support is not installed on this server.") from exc

    prs = Presentation()

    def _brand(slide):
        """The school's name and logo along the foot of a slide.

        A deck shown at a parents' meeting or to the board carries the school's
        identity on every slide, not only the first, for the same reason the PDF
        letterhead repeats: slides get photographed and shared one at a time.
        """
        if not letterhead:
            return
        try:
            if os.path.exists(LOGO_PATH):
                slide.shapes.add_picture(
                    LOGO_PATH, Inches(0.25), prs.slide_height - Inches(0.62), height=Inches(0.45)
                )
            box = slide.shapes.add_textbox(
                Inches(0.85), prs.slide_height - Inches(0.55),
                prs.slide_width - Inches(1.2), Inches(0.4),
            )
            frame = box.text_frame
            frame.text = " | ".join(letterhead_lines()[:2])
            run = frame.paragraphs[0].runs[0]
            run.font.size = Pt(8)
        except Exception:  # pragma: no cover - branding must never cost the deck
            logger.warning("PPTX: the slide branding could not be placed")

    if title:
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = str(title)[:300]
        if len(cover.placeholders) > 1 and paragraphs:
            cover.placeholders[1].text = _clean_cell(paragraphs[0])
        _brand(cover)

    body_slides = slides or []
    if not body_slides and (paragraphs or rows):
        body_slides = [{"title": "Details", "bullets": [_clean_cell(p) for p in (paragraphs or [])[1:]]}]

    for spec in body_slides[:MAX_SLIDES]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _clean_cell(spec.get("title", ""))[:200]
        bullets = spec.get("bullets") or []
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for i, bullet in enumerate(bullets[:20]):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = _clean_cell(bullet)
            para.font.size = Pt(18)
        _brand(slide)

    if rows:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Data"
        # A slide cannot hold thousands of rows legibly; show a window and say so.
        shown = rows[:12]
        table_rows = len(shown) + (1 if headers else 0)
        shape = slide.shapes.add_table(
            table_rows, len(shown[0]), Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * table_rows)
        )
        table = shape.table
        offset = 0
        if headers:
            for i, h in enumerate(headers):
                table.cell(0, i).text = h
            offset = 1
        for r, row in enumerate(shown):
            for c, value in enumerate(row):
                table.cell(r + offset, c).text = value
        _brand(slide)

    if truncated_note:
        note_slide = prs.slides.add_slide(prs.slide_layouts[5])
        note_slide.shapes.title.text = truncated_note[:200]
        _brand(note_slide)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pdf(title, paragraphs, headers, rows, truncated_note, letterhead=True):
    try:
        from fpdf import FPDF
    except ImportError as exc:  # pragma: no cover
        raise DocumentBuildError("PDF support is not installed on this server.") from exc

    letterhead_cls = _letterhead_pdf_class(FPDF, _identity()) if letterhead else FPDF
    pdf = letterhead_cls(orientation="P", unit="mm", format="A4")

    # Fonts are registered BEFORE the first page, because `add_page()` immediately
    # calls `header()`, which draws in whichever font is set at that moment. Adding
    # the page first, as this did originally, meant the letterhead was drawn before
    # the Unicode fonts existed.
    unicode_ok = _register_unicode_font(pdf)
    enc = _build_pdf_text_encoder(unicode_ok)
    family = UNICODE_FONT_FAMILY if unicode_ok else "Helvetica"
    pdf._family = family
    pdf._enc = enc

    # Both Noto families ship Regular and Bold only. Asking fpdf2 for italic would
    # raise and lose the document, so the closing note is set in regular weight when
    # the Unicode fonts are in use. A missing italic is a far smaller loss than a
    # missing document.
    note_style = "" if unicode_ok else "I"

    # Room for the letterhead at the top and the address block at the foot, so body
    # text never runs into either.
    top_margin = 34 if letterhead else 15
    pdf.set_margins(left=15, top=top_margin, right=15)
    pdf.set_auto_page_break(auto=True, margin=20 if letterhead else 15)
    pdf.alias_nb_pages()  # makes "Page 2 of 7" resolve, rather than printing "{nb}"
    pdf.add_page()

    if title:
        pdf.set_font(family, "B", 15)
        pdf.set_text_color(*BRAND_NAVY)
        pdf.multi_cell(0, 8, enc(str(title)[:300]), align="C")
        pdf.set_text_color(0, 0, 0)
        pdf.ln(3)

    pdf.set_font(family, size=10.5)
    for para in (paragraphs or [])[:MAX_PARAGRAPHS]:
        pdf.multi_cell(0, 5.6, enc(_clean_cell(para)))
        pdf.ln(1.5)

    if rows:
        pdf.ln(3)
        usable = pdf.w - pdf.l_margin - pdf.r_margin
        col_count = len(rows[0]) or 1
        col_width = usable / col_count

        def _draw_header_row():
            pdf.set_font(family, "B", 8.5)
            pdf.set_fill_color(*BRAND_NAVY)
            pdf.set_text_color(255, 255, 255)
            pdf.set_draw_color(*RULE_GREY)
            for h in headers:
                pdf.cell(col_width, 7, enc(h)[:40], border=1, fill=True, align="C")
            pdf.ln()
            pdf.set_text_color(0, 0, 0)

        if headers:
            _draw_header_row()

        pdf.set_font(family, size=8.5)
        pdf.set_draw_color(*RULE_GREY)
        for index, row in enumerate(rows):
            if pdf.get_y() > pdf.h - 28:
                pdf.add_page()
                # The heading row is repeated on each new page. A table continuing
                # over a page break with no headings is unreadable on paper, which is
                # how most of these are actually used.
                if headers:
                    _draw_header_row()
                    pdf.set_font(family, size=8.5)
            # Banded rows: on a wide fee sheet, following one student across eight
            # columns without them is where the eye slips a line.
            shaded = index % 2 == 1
            if shaded:
                pdf.set_fill_color(243, 246, 251)
            for value in row:
                pdf.cell(col_width, 6, enc(value)[:40], border=1, fill=shaded)
            pdf.ln()

    if truncated_note:
        pdf.ln(4)
        pdf.set_font(family, note_style, 8.5)
        pdf.set_text_color(90, 100, 120)
        pdf.multi_cell(0, 5, enc(truncated_note))
        pdf.set_text_color(0, 0, 0)

    out = pdf.output(dest="S")
    return bytes(out) if isinstance(out, (bytes, bytearray)) else out.encode("latin-1")


def _build_xml(title, paragraphs, headers, rows, truncated_note):
    """Build a well-formed XML document from the supplied content.

    The structure is intentionally simple so any XML parser can read it without
    a schema: a root <document> with <metadata>, <content> (prose paragraphs),
    and <table> (headers + rows) sections. The school's name is in metadata, not
    embedded as extra rows, so a program iterating <row> elements gets clean data.
    """
    import xml.etree.ElementTree as ET

    root = ET.Element("document")
    if title:
        root.set("title", str(title)[:300])

    # Metadata carries the school's identity without polluting the data rows.
    meta = ET.SubElement(root, "metadata")
    try:
        from school_identity import default_school_identity
        identity = default_school_identity()
        ET.SubElement(meta, "school").text = identity.get("school_name", "")
        ET.SubElement(meta, "board").text = identity.get("board", "")
        ET.SubElement(meta, "city").text = identity.get("city", "")
    except Exception:
        pass
    if title:
        ET.SubElement(meta, "title").text = str(title)[:300]
    if truncated_note:
        ET.SubElement(meta, "note").text = truncated_note

    if paragraphs:
        content_el = ET.SubElement(root, "content")
        for para in paragraphs[:MAX_PARAGRAPHS]:
            p = ET.SubElement(content_el, "paragraph")
            p.text = _clean_cell(para)

    if rows:
        table_el = ET.SubElement(root, "table")
        if headers:
            hdrs_el = ET.SubElement(table_el, "headers")
            for h in headers:
                ET.SubElement(hdrs_el, "header").text = h
        rows_el = ET.SubElement(table_el, "rows")
        for row in rows:
            row_el = ET.SubElement(rows_el, "row")
            for idx, cell in enumerate(row):
                cell_el = ET.SubElement(row_el, "cell")
                if headers and idx < len(headers):
                    # Use the column name as the attribute so the cell is self-describing.
                    safe_attr = re.sub(r"[^A-Za-z0-9_.-]", "_", headers[idx]) or f"col{idx}"
                    cell_el.set("name", safe_attr)
                cell_el.text = cell

    # ET.tostring gives bytes with the declaration only when encoding != "unicode".
    declaration = b'<?xml version="1.0" encoding="UTF-8"?>\n'
    body = ET.tostring(root, encoding="unicode", xml_declaration=False)
    return declaration + body.encode("utf-8")


def _build_text(doc_type, title, paragraphs, headers, rows, truncated_note, letterhead=True):
    lines: List[str] = []
    if doc_type == "csv":
        import csv

        buf = io.StringIO()
        writer = csv.writer(buf)
        if letterhead:
            # A CSV has no header region, so the branding has to become rows.
            #
            # NOTE THE TRADE-OFF, because it is a real one: another system reading
            # this file will now see three rows before the column names. That is the
            # cost of branding a data-interchange format, and it is why `build_document`
            # takes `letterhead=False` - anything being fed to another program should
            # pass it.
            for line in letterhead_lines():
                writer.writerow([line])
            writer.writerow([])
        if headers:
            writer.writerow(headers)
        writer.writerows(rows)
        if truncated_note:
            writer.writerow([truncated_note])
        if letterhead:
            writer.writerow([])
            for line in letterhead_footer_lines():
                writer.writerow([line])
        return buf.getvalue().encode("utf-8")

    if letterhead:
        head = letterhead_lines()
        if head:
            if doc_type == "md":
                lines.append(f"**{head[0]}**")
                lines.extend(head[1:])
                lines.append("")
                lines.append("---")
            else:
                lines.extend(head)
                lines.append("=" * max(len(line) for line in head))
            lines.append("")

    if title:
        lines.append(f"# {title}" if doc_type == "md" else str(title))
        lines.append("")
    lines.extend(_clean_cell(p) for p in (paragraphs or [])[:MAX_PARAGRAPHS])
    if rows:
        lines.append("")
        if doc_type == "md" and headers:
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("|" + "|".join([" --- "] * len(headers)) + "|")
            lines.extend("| " + " | ".join(r) + " |" for r in rows)
        else:
            if headers:
                lines.append("\t".join(headers))
            lines.extend("\t".join(r) for r in rows)
    if truncated_note:
        lines.extend(["", truncated_note])
    if letterhead:
        foot = letterhead_footer_lines()
        if foot:
            lines.append("")
            lines.append("---" if doc_type == "md" else "-" * 60)
            lines.extend(foot)
    return "\n".join(lines).encode("utf-8")


# ── The one entry point ─────────────────────────────────────────────────────────

def _workbook_cell(value: Any):
    """A cell value Excel will treat as what it is.

    NUMBERS STAY NUMBERS. Everywhere else in this module a cell is stringified, which
    is harmless in a letter and wrong in a workbook: the office downloads the fee sheet
    in order to ADD IT UP, and "12400" as text will not add. The same rule was already
    settled for the hand-wired downloads in Release 3 item 6; this is that rule where
    the whole school's money lands.

    Booleans are deliberately NOT passed through as numbers - Excel would show TRUE as
    1 next to a column of rupees, which reads as an amount.
    """
    if value is None:
        return ""
    if isinstance(value, bool):
        return "Yes" if value else "No"
    if isinstance(value, (int, float)):
        return value
    return _clean_cell(value)


def build_workbook(
    *,
    sheets: List[Dict[str, Any]],
    filename: str = "",
    max_rows: Optional[int] = None,
) -> BuiltDocument:
    """One Excel file holding several sheets, one per area of the school.

    `sheets` is a list of `{"name": str, "headers": [...], "rows": [[...]]}`.

    WHY THIS IS NOT `build_document` CALLED SEVEN TIMES. openpyxl can only write one
    workbook per save, so a per-sheet loop would produce seven separate files. The
    whole point of the whole-school download is that the office opens ONE file and
    tabs across it.

    NOTHING IS EVER TRIMMED HERE. `max_rows` is a REFUSAL ceiling: a sheet over it
    raises rather than dropping rows, because this file leaves the building and gets
    filed as the school's record. A short sheet inside a workbook is the worst version
    of this release's defining fault - nobody scrolls to the bottom of tab five to
    check whether it stopped early.

    Each sheet says its own row count on its first line, above the headings, for the
    same reason: a person cannot count 1,876 rows by eye, so the file has to tell them.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except ImportError as exc:  # pragma: no cover
        raise DocumentBuildError("Excel support is not installed on this server.") from exc

    if not sheets:
        raise DocumentBuildError("A workbook needs at least one sheet.")

    ceiling = MAX_ROWS if max_rows is None else max_rows

    wb = Workbook()
    wb.remove(wb.active)
    used: set = set()
    for sheet in sheets:
        raw_name = str(sheet.get("name") or "Sheet")
        # Excel refuses to OPEN a file with an illegal or duplicated sheet name rather
        # than complaining about it, so this is silently fatal if not handled.
        name = (re.sub(r"[\[\]:*?/\\]", "-", raw_name)[:31]) or "Sheet"
        suffix = 2
        while name.lower() in used:
            tail = f" {suffix}"
            name = name[: 31 - len(tail)] + tail
            suffix += 1
        used.add(name.lower())

        headers = [_clean_cell(h) for h in (sheet.get("headers") or [])][:MAX_COLUMNS]
        rows = sheet.get("rows") or []
        if len(rows) > ceiling:
            raise DocumentBuildError(
                f"The '{raw_name}' sheet has {len(rows):,} rows, which is more than the "
                f"{ceiling:,} one file may hold. No file was produced, and nothing was "
                "left out of one: narrow it with a date range and download in parts."
            )

        ws = wb.create_sheet(title=name)
        ws.append([f"{raw_name}: {len(rows):,} rows"])
        ws.cell(row=1, column=1).font = Font(bold=True)
        ws.append([])
        if headers:
            ws.append(headers)
            for cell in ws[ws.max_row]:
                cell.font = Font(bold=True)
        for row in rows:
            if not isinstance(row, (list, tuple)):
                row = [row]
            ws.append([_workbook_cell(c) for c in row][:MAX_COLUMNS])

        for idx, column_cells in enumerate(ws.columns, start=1):
            longest = max((len(str(c.value)) for c in column_cells if c.value is not None), default=0)
            ws.column_dimensions[ws.cell(row=1, column=idx).column_letter].width = min(60, max(10, longest + 2))

    buf = io.BytesIO()
    wb.save(buf)
    return BuiltDocument(
        content=buf.getvalue(),
        content_type=CONTENT_TYPES["xlsx"],
        filename=safe_filename(filename or "whole-school", "xlsx"),
        doc_type="xlsx",
        truncated=False,
    )


def build_document(
    *,
    doc_type: str,
    filename: str = "",
    title: str = "",
    paragraphs: Optional[List[Any]] = None,
    headers: Optional[List[Any]] = None,
    rows: Optional[List[List[Any]]] = None,
    slides: Optional[List[Dict[str, Any]]] = None,
    letterhead: bool = True,
    max_rows: Optional[int] = None,
) -> BuiltDocument:
    """Build a document and return its bytes. Never touches the database or S3.

    `letterhead` puts the school's own branding on the document. On by default:
    nearly everything generated here is sent to a parent or filed by the school, and a
    page with no letterhead does not read as a school document. Pass False for
    something purely internal.

    EXCEPT for spreadsheets - see `UNBRANDED_TYPES` below.

    `max_rows` is the row ceiling for the table. Left unset it is MAX_ROWS (5,000),
    which is a TRUNCATION point: rows past it are dropped and the file carries a note
    saying so. That trade is right for Flo's generated documents, where a reader is
    looking at a summary. **It is wrong for an export**, which leaves the building and
    gets filed as a record - the fee ledger alone is about 10,700 rows, so an Excel
    download of it was quietly losing more than half. `routes/exports.py` therefore
    passes its own refusal ceiling here, so the only limit an export can hit is the
    one that refuses the request outright rather than shipping a short file.


    Raises DocumentBuildError for anything malformed, before any caller has stored
    something it would then have to clean up.
    """
    doc_type = (doc_type or "").lower().lstrip(".")
    # Spreadsheets are DATA, not stationery (Abhimanyu, 2026-08-07). Branding rows
    # above the column headings push every row down and mean anything reading the
    # file - a formula, a filter, another program, an import into another system -
    # starts on the wrong line. The gain was cosmetic and the cost was real, so the
    # branding is off for these two regardless of what the caller asks for.
    if doc_type in UNBRANDED_TYPES:
        letterhead = False
    _validate(doc_type, title)

    if not any([title, paragraphs, rows, slides]):
        raise DocumentBuildError("Nothing to put in the document - provide a title, text, rows or slides.")

    hdrs, norm_rows, truncated = _normalise_table(headers, rows, max_rows)
    note = ""
    if truncated:
        # Say it in the file itself. A silently short export is the Epic 4 defect
        # (a failure that looks like a complete answer) in a new place.
        applied = MAX_ROWS if max_rows is None else max_rows
        note = f"Note: only the first {applied:,} rows are included. {len(rows):,} rows matched."

    if doc_type == "docx":
        content = _build_docx(title, paragraphs, hdrs, norm_rows, note, letterhead)
    elif doc_type == "xlsx":
        content = _build_xlsx(title, paragraphs, hdrs, norm_rows, note, letterhead)
    elif doc_type == "pptx":
        content = _build_pptx(title, paragraphs, hdrs, norm_rows, note, slides, letterhead)
    elif doc_type == "pdf":
        content = _build_pdf(title, paragraphs, hdrs, norm_rows, note, letterhead)
    elif doc_type == "xml":
        content = _build_xml(title, paragraphs, hdrs, norm_rows, note)
    else:
        content = _build_text(doc_type, title, paragraphs, hdrs, norm_rows, note, letterhead)

    return BuiltDocument(
        content=content,
        content_type=CONTENT_TYPES[doc_type],
        filename=safe_filename(filename or title or "document", doc_type),
        doc_type=doc_type,
        truncated=truncated,
        notes=[note] if note else [],
        editable_html=editable_html(title, paragraphs, hdrs, norm_rows, note),
    )
