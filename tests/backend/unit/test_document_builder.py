"""UI Sweep Epic 10, Story 10.1 - the shared document builder.

These assert the bytes are a REAL file of the claimed type, by opening them again with
the corresponding reader. Asserting "no exception was raised" would pass for a builder
that wrote an empty or corrupt file, which is exactly the class of defect that would
reach the owner as "the download won't open".
"""
from __future__ import annotations

import io
import os
import zipfile

import pytest

from services import document_builder
from services.document_builder import (
    CONTENT_TYPES,
    SUPPORTED_TYPES,
    UNBRANDED_TYPES,
    MAX_ROWS,
    DocumentBuildError,
    build_document,
    safe_filename,
)

# ── The files are real files ────────────────────────────────────────────────────

def test_xlsx_opens_as_a_workbook_with_the_data_in_it():
    doc = build_document(
        doc_type="xlsx",
        title="Fee Sheet",
        headers=["Student", "Class", "Owed"],
        rows=[["Asha", "5-A", "12000"], ["Bipin", "3-B", "9000"]],
    )
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(doc.content))
    ws = wb.active
    values = [[c.value for c in row] for row in ws.iter_rows()]
    assert ["Student", "Class", "Owed"] in values
    assert ["Asha", "5-A", "12000"] in values
    assert doc.content_type == CONTENT_TYPES["xlsx"]


def test_docx_opens_as_a_document_with_the_text_in_it():
    doc = build_document(
        doc_type="docx",
        title="Principal's Circular",
        paragraphs=["School reopens on 1 April.", "Uniforms are compulsory."],
    )
    from docx import Document

    parsed = Document(io.BytesIO(doc.content))
    text = "\n".join(p.text for p in parsed.paragraphs)
    assert "Principal's Circular" in text
    assert "School reopens on 1 April." in text


def test_pptx_opens_as_a_presentation():
    doc = build_document(
        doc_type="pptx",
        title="School Profile",
        slides=[{"title": "Results", "bullets": ["100% pass", "12 distinctions"]}],
    )
    from pptx import Presentation

    prs = Presentation(io.BytesIO(doc.content))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "School Profile" in titles
    assert "Results" in titles


def test_pdf_is_a_pdf():
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    assert doc.content.startswith(b"%PDF")
    assert doc.content_type == "application/pdf"


def test_office_files_are_valid_zip_containers():
    """A .docx/.xlsx/.pptx is a zip. A truncated write produces bytes that look
    plausible and fail to open - this catches that without a full parse."""
    for doc_type in ("docx", "xlsx", "pptx"):
        doc = build_document(doc_type=doc_type, title="T", paragraphs=["body"])
        assert zipfile.is_zipfile(io.BytesIO(doc.content)), doc_type


def test_csv_and_markdown_come_out_as_text():
    # CSV is in UNBRANDED_TYPES, so it is plain either way. Markdown passes
    # `letterhead=False` because this test is about the DATA shape; the branded form
    # is covered by `test_every_readable_format_carries_the_schools_branding`.
    csv_doc = build_document(doc_type="csv", headers=["A", "B"], rows=[["1", "2"]])
    assert csv_doc.content.decode().splitlines()[0] == "A,B"

    md_doc = build_document(doc_type="md", title="Policy", headers=["A"], rows=[["1"]], letterhead=False)
    text = md_doc.content.decode()
    assert text.startswith("# Policy")
    assert "| A |" in text


# ── Filenames are the security-sensitive part ───────────────────────────────────

@pytest.mark.parametrize("dangerous", [
    "../../../etc/passwd",
    "..\\..\\windows\\system32\\config",
    'fee"; rm -rf /',
    "report\nContent-Disposition: attachment; filename=evil",
    "  ",
    "",
])
def test_filenames_cannot_escape_or_forge_a_header(dangerous):
    """The name may come from a person or from Flo. A path separator walks outside
    the school's S3 prefix; a newline or quote forges a response header."""
    name = safe_filename(dangerous, "xlsx")
    assert "/" not in name and "\\" not in name
    assert "\n" not in name and "\r" not in name
    assert '"' not in name and ";" not in name
    assert ".." not in name
    assert name.endswith(".xlsx")
    assert len(name) > len(".xlsx")


def test_a_traversal_attempt_keeps_only_the_last_component():
    assert safe_filename("../../secret-report", "pdf") == "secret-report.pdf"


def test_the_extension_is_not_doubled():
    assert safe_filename("fees.xlsx", "xlsx") == "fees.xlsx"


def test_a_very_long_name_is_cut():
    assert len(safe_filename("x" * 500, "pdf")) <= 85


# ── Refusing, rather than writing something broken ──────────────────────────────

def test_an_unsupported_type_is_refused():
    with pytest.raises(DocumentBuildError) as exc:
        build_document(doc_type="exe", title="x")
    assert "Unsupported" in str(exc.value)


def test_an_empty_description_is_refused():
    """Refused BEFORE storage, so no orphan object is left in S3 with no record."""
    with pytest.raises(DocumentBuildError):
        build_document(doc_type="docx")


def test_an_absurd_title_is_refused():
    with pytest.raises(DocumentBuildError):
        build_document(doc_type="docx", title="t" * 400)


# ── Honest truncation ───────────────────────────────────────────────────────────

def test_too_many_rows_are_cut_and_the_file_says_so():
    """A silently short export is the Epic 4 defect - a failure that looks like a
    complete answer - in a new place."""
    rows = [[str(i), "x"] for i in range(MAX_ROWS + 500)]
    doc = build_document(doc_type="csv", headers=["N", "V"], rows=rows)

    assert doc.truncated is True
    text = doc.content.decode()
    assert "only the first" in text
    assert f"{MAX_ROWS + 500:,}" in text


def test_a_normal_export_is_not_marked_truncated():
    doc = build_document(doc_type="csv", headers=["N"], rows=[["1"], ["2"]])
    assert doc.truncated is False
    assert doc.notes == []


# ── Real data is ragged ─────────────────────────────────────────────────────────

def test_ragged_rows_are_padded_rather_than_rejected():
    """Refusing a whole export because one student has no phone number would be
    worse than filling a blank."""
    doc = build_document(
        doc_type="csv",
        headers=["Name", "Class", "Phone"],
        rows=[["Asha", "5-A", "99999"], ["Bipin"], ["Chetan", "3-B"]],
    )
    lines = doc.content.decode().strip().splitlines()
    assert all(line.count(",") == 2 for line in lines), lines


def test_none_values_become_blanks_not_the_word_none():
    doc = build_document(doc_type="csv", headers=["A", "B"], rows=[[None, "x"]])
    assert "None" not in doc.content.decode()


def test_devanagari_does_not_lose_the_whole_pdf():
    """A Hindi circular must produce a file rather than raising and losing everything
    the user asked for. This held even in the Latin-1 days, and must keep holding."""
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["आज छुट्टी है"])
    assert doc.content.startswith(b"%PDF")


# ── Hindi in PDFs (2026-08-07) ──────────────────────────────────────────────────
#
# The whole point of the change: Devanagari used to be replaced with "?" silently, so
# a circular to parents in Amroha came out as a row of question marks and the download
# still looked like it had worked. These tests assert the text is ACTUALLY IN THE FILE,
# by reading the PDF back. Asserting "no exception was raised" would have passed
# happily throughout the entire period the defect existed.

def _pdf_text(content: bytes) -> str:
    from pypdf import PdfReader

    return PdfReader(io.BytesIO(content)).pages[0].extract_text()


def test_hindi_reaches_the_page_instead_of_becoming_question_marks():
    doc = build_document(
        doc_type="pdf",
        title="विद्यालय सूचना",
        paragraphs=["आज विद्यालय बंद रहेगा"],
    )
    text = _pdf_text(doc.content)
    # The specific failure being guarded: every Devanagari character replaced by "?".
    assert "?" not in text
    # Devanagari is U+0900–U+097F. At least some of it must have reached the page.
    assert any("ऀ" <= ch <= "ॿ" for ch in text), text


def test_hindi_in_a_table_reaches_the_page_too():
    """Row and header cells go through a different call than paragraphs do, so they
    are checked separately - a fix applied to only one of the two is the likelier
    mistake than no fix at all."""
    doc = build_document(
        doc_type="pdf",
        title="उपस्थिति",
        headers=["नाम", "कक्षा"],
        rows=[["अजमल", "III"]],
    )
    text = _pdf_text(doc.content)
    assert "?" not in text
    assert any("ऀ" <= ch <= "ॿ" for ch in text), text


def test_english_documents_are_unaffected():
    """THIS TEST CAUGHT A REAL REGRESSION, so do not weaken it.

    The first attempt registered Noto Sans Devanagari as the only font, assuming it
    covered Latin too. It does not. This came back as '  .' - every letter of
    "Holiday on Monday." silently dropped, leaving the space and the full stop. Almost
    every document this school generates is in English, so that would have been a far
    worse defect than the Hindi one being fixed."""
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    text = _pdf_text(doc.content)
    assert "Holiday on Monday." in text


def test_a_sentence_mixing_hindi_and_english_keeps_both():
    """The realistic case, and the one neither font handles alone: a class number in
    Latin digits inside a Hindi sentence. It works because the Latin font is the
    document font and the Devanagari font is registered as its per-character
    fallback."""
    doc = build_document(
        doc_type="pdf",
        title="Notice",
        paragraphs=["कक्षा III की उपस्थिति 92% रही।"],
    )
    text = _pdf_text(doc.content)
    assert "III" in text, f"the Latin part was dropped: {text!r}"
    assert any("ऀ" <= ch <= "ॿ" for ch in text), f"the Hindi part was dropped: {text!r}"


def test_a_missing_font_degrades_to_the_old_behaviour_rather_than_failing(monkeypatch):
    """The font file not being on the server must never take document generation down
    for the whole school. It falls back to the pre-2026-08-07 Latin-1 substitution,
    which is lossy but is a working document, and it says so in the log."""
    monkeypatch.setattr(document_builder, "unicode_font_available", lambda: False)
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["आज छुट्टी है"])
    assert doc.content.startswith(b"%PDF")


def test_the_shipped_configuration_has_both_the_font_and_the_shaper():
    """Hindi is only rendered CORRECTLY with both. Without the shaper the characters
    are present but conjuncts and vowel signs sit in the wrong places, which is a
    different and quieter kind of wrong. If this fails, the deploy is missing a piece."""
    assert document_builder.unicode_font_available(), (
        f"Devanagari font missing from {document_builder.FONT_DIR}"
    )
    assert document_builder.text_shaping_available(), "uharfbuzz is not installed"


# ── The school's letterhead (2026-08-07) ────────────────────────────────────────

def test_the_letterhead_matches_the_schools_printed_form_word_for_word():
    """Abhimanyu compared a generated PDF against the school's printed enquiry form on
    2026-08-07 and found three differences: the name was set as "The Aaryans" where the
    paper says "THE AARYANS", the footer had been reworded out of the identity fields
    rather than transcribed, and the pale wordmark was missing.

    So the letterhead now takes its strings from `school_identity.LETTERHEAD`, which is
    the printed wording character for character. Rebuilding these out of `address`,
    `phone` and `email` is exactly what produced the mismatch, so this test asserts the
    printed forms and not the identity fields."""
    from school_identity import LETTERHEAD

    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    text = _pdf_text(doc.content)

    assert LETTERHEAD["name"] in text
    assert LETTERHEAD["tagline"] in text
    assert LETTERHEAD["affiliation_line"] in text
    # The footer lines are the ones that had been reworded, so they are checked whole.
    assert LETTERHEAD["footer_address"] in text
    assert LETTERHEAD["footer_contact"] in text


def test_the_pale_wordmark_is_printed_across_the_page():
    """The printed form carries "THE AARYANS" tiled faintly across the whole sheet,
    behind the text. A single occurrence would just be the heading, so this counts:
    the watermark has to actually repeat."""
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    text = _pdf_text(doc.content)
    assert text.count("THE AARYANS") > 20, (
        "the tiled wordmark is missing or barely drawn; the first version of it was "
        "lost entirely because the tiling started at a negative x, which fpdf2 refuses"
    )


def _watermark_positions(content: bytes):
    """Where each piece of watermark text was drawn, in PDF points from the top left.

    Uses pypdf, which is already a dependency, rather than PyMuPDF. PyMuPDF would give
    this more directly but it is AGPL, and pulling an AGPL library into the backend to
    satisfy one test is not a trade worth making.
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    page = reader.pages[0]
    height = float(page.mediabox.height)
    found = []

    def visit(text, cm, tm, font_dict, font_size):
        if "AARYANS" not in (text or "").upper():
            return
        # The letterhead's own name and its footer lines also contain the word, so
        # they are excluded by type size. Without this the row count came out three
        # too high, which is the heading plus the two footer lines.
        if font_size is None or abs(font_size) > document_builder.WATERMARK_FONT_PT + 1:
            return
        # tm[4]/tm[5] are the drawing position. PDF measures y upward from the
        # bottom, so it is flipped to match the builder's top-down thinking.
        found.append((tm[4], height - tm[5]))

    page.extract_text(visitor_text=visit)
    return found


def test_the_wordmark_breaks_around_the_crest_instead_of_printing_over_it():
    """On the printed form the repeated wording stops at the crest's edge and picks up
    again on the other side. The first version tiled straight across it and turned the
    crest into a smudge (Abhimanyu, 2026-08-07, with a screenshot).

    Checked by geometry, not by eye: nothing may be drawn inside the crest's box."""
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])

    pt = 72 / 25.4  # the builder works in millimetres; the PDF's unit is the point
    size = document_builder.CHARIOT_SIZE_MM * pt
    from pypdf import PdfReader

    page = PdfReader(io.BytesIO(doc.content)).pages[0]
    page_w, page_h = float(page.mediabox.width), float(page.mediabox.height)
    left, top = (page_w - size) / 2, (page_h - size) / 2
    right, bottom = left + size, top + size

    intruders = [
        (x, y) for x, y in _watermark_positions(doc.content)
        if left < x < right and top < y < bottom
    ]
    assert not intruders, f"watermark text printed over the crest: {intruders[:3]}"


def test_the_wordmark_is_laid_out_in_the_rows_counted_off_the_printed_form():
    """Fourteen rows down the sheet with a small gap, counted off the printed form."""
    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    rows = {round(y) for _, y in _watermark_positions(doc.content)}
    # Rounding can split one row across two neighbouring values, so a little slack
    # rather than demanding exactly WATERMARK_ROWS.
    assert abs(len(rows) - document_builder.WATERMARK_ROWS) <= 2, sorted(rows)


def test_the_background_crest_is_the_pale_copy_not_the_full_colour_logo():
    """Full colour behind a fee table reads as a mistake; the printed form shows a
    pale ghost. Two files, so the letterhead crest at the top stays full colour."""
    assert os.path.exists(document_builder.CREST_WATERMARK_PATH), (
        "the desaturated background crest is missing, so the watermark would fall "
        "back to the full-colour logo"
    )


def test_the_watermark_never_creates_a_page_of_its_own():
    """The tiling loop writes near the bottom of the sheet. With automatic page breaks
    left on, that tripped a page break, which called the header again, which drew
    another watermark, until Python ran out of stack. A one-paragraph notice must be
    exactly one page."""
    from pypdf import PdfReader

    doc = build_document(doc_type="pdf", title="Notice", paragraphs=["Holiday on Monday."])
    assert len(PdfReader(io.BytesIO(doc.content)).pages) == 1


def test_the_letterhead_is_on_every_page_not_just_the_first():
    """A twelve-page fee report that stops looking like a school document after page
    one is the failure this guards. fpdf2 calls header()/footer() per page, which is
    why the letterhead is a subclass rather than something drawn once."""
    rows = [[f"Student {i}", f"{i}-A", "94%"] for i in range(1, 60)]
    doc = build_document(
        doc_type="pdf", title="Fee Report", headers=["Name", "Class", "Attendance"], rows=rows
    )
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(doc.content))
    assert len(reader.pages) > 1, "the fixture no longer spans pages, so it proves nothing"
    for number, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        assert "THE AARYANS" in text, f"page {number} has no letterhead"
        assert f"Page {number} of {len(reader.pages)}" in text, f"page {number} is not numbered"


def test_a_table_repeats_its_headings_after_a_page_break():
    """On paper, a table continuing onto a second page with no column headings is
    unreadable, and paper is how the school actually uses these."""
    rows = [[f"Student {i}", f"{i}-A", "94%"] for i in range(1, 60)]
    doc = build_document(
        doc_type="pdf", title="Fee Report", headers=["Name", "Class", "Attendance"], rows=rows
    )
    from pypdf import PdfReader

    second = PdfReader(io.BytesIO(doc.content)).pages[1].extract_text()
    assert "Attendance" in second, "the column headings were not repeated"


def test_every_readable_format_carries_the_schools_branding():
    """Abhimanyu, 2026-08-07: "make sure that this template is used for each and every
    type of document as being the template and branding of the school", then, once he
    had seen it: "remove the branding from excel/csv file type".

    So: every format a PERSON reads is branded. Spreadsheets are data and are not.

    Derived from `SUPPORTED_TYPES` minus `UNBRANDED_TYPES` rather than a hand-written
    list, so a format added later fails this test until it is branded too. That is the
    point: the first version branded PDF and Word and quietly left the rest plain."""
    from school_identity import LETTERHEAD

    name = LETTERHEAD["name"]
    for doc_type in sorted(set(SUPPORTED_TYPES) - UNBRANDED_TYPES):
        doc = build_document(
            doc_type=doc_type,
            title="Notice",
            paragraphs=["Holiday on Monday."],
            headers=["A"],
            rows=[["1"]],
        )
        if doc_type == "pdf":
            haystack = _pdf_text(doc.content)
        elif doc_type in ("docx", "xlsx", "pptx"):
            # The Office formats are zip containers; the branding lands in different
            # parts of each (Word's header part, Excel's sheet, PowerPoint's slides),
            # so the container as a whole is searched rather than one known part.
            with zipfile.ZipFile(io.BytesIO(doc.content)) as zf:
                haystack = "".join(
                    zf.read(entry).decode("utf-8", "ignore")
                    for entry in zf.namelist()
                    if entry.endswith(".xml")
                )
        else:
            haystack = doc.content.decode("utf-8")
        assert name in haystack, f"{doc_type} carries no school branding"


def test_branding_can_be_turned_off_for_every_format():
    """The switch has to work everywhere, not only on the two formats it was first
    written for."""
    from school_identity import LETTERHEAD

    for doc_type in ("csv", "md", "txt"):
        doc = build_document(
            doc_type=doc_type, title="Data", headers=["A"], rows=[["1"]], letterhead=False
        )
        assert LETTERHEAD["name"] not in doc.content.decode("utf-8"), doc_type


def test_a_spreadsheet_is_never_branded_even_when_asked_for():
    """A spreadsheet is DATA. Branding rows above the column headings shift every row
    down, so a formula, a filter or an import into another system starts on the wrong
    line. The gain was cosmetic and the cost was real, so `letterhead=True` is ignored
    for these two rather than merely defaulting to off (Abhimanyu, 2026-08-07)."""
    from school_identity import LETTERHEAD
    from openpyxl import load_workbook

    csv_doc = build_document(
        doc_type="csv", headers=["A", "B"], rows=[["1", "2"]], letterhead=True
    )
    text = csv_doc.content.decode()
    assert LETTERHEAD["name"] not in text
    # The very first line is the column headings, with nothing above them.
    assert text.splitlines()[0] == "A,B"

    xlsx_doc = build_document(
        doc_type="xlsx", title="Fees", headers=["A", "B"], rows=[["1", "2"]], letterhead=True
    )
    ws = load_workbook(io.BytesIO(xlsx_doc.content)).active
    values = [[c.value for c in row] for row in ws.iter_rows()]
    assert values[0] == ["A", "B"], values[:3]
    assert all(LETTERHEAD["name"] not in str(cell) for row in values for cell in row)
    # And not hidden in the print header or footer either.
    assert LETTERHEAD["name"] not in str(ws.oddHeader.center.text or "")
    assert LETTERHEAD["name"] not in str(ws.oddFooter.center.text or "")


def test_the_letterhead_can_be_turned_off():
    doc = build_document(
        doc_type="pdf", title="Internal", paragraphs=["Scratch note."], letterhead=False
    )
    assert "Affiliation No." not in _pdf_text(doc.content)


def test_word_documents_get_the_letterhead_in_the_page_furniture():
    """Written into the Word section header, not as body paragraphs, so it repeats on
    every page and cannot be half-deleted by someone editing the circular."""
    from docx import Document

    doc = build_document(doc_type="docx", title="Circular", paragraphs=["Holiday on Monday."])
    parsed = Document(io.BytesIO(doc.content))
    header_text = "\n".join(p.text for p in parsed.sections[0].header.paragraphs)
    footer_text = "\n".join(p.text for p in parsed.sections[0].footer.paragraphs)
    assert "THE AARYANS" in header_text
    assert "Delhi-Moradabad Highway" in footer_text


def test_the_logo_ships_with_the_backend():
    """The letterhead reads the logo off disk at build time. If the file is not in the
    deploy the header quietly loses its logo, which nobody would notice from a test
    that only checked the text."""
    assert os.path.exists(document_builder.LOGO_PATH), (
        f"the school logo is missing from {document_builder.LOGO_PATH}"
    )


def test_the_font_licence_ships_beside_the_font():
    """SIL OFL 1.1 permits redistribution and embedding on condition the licence
    travels with the font. `.ebignore` excludes every *.txt, so this file is only in
    the deploy because of an explicit exception there - easy to lose by accident."""
    licence = os.path.join(document_builder.FONT_DIR, "OFL.txt")
    assert os.path.exists(licence), "the font's licence file is missing"
    assert "SIL Open Font License" in open(licence, encoding="utf-8").read()


def test_devanagari_survives_intact_in_the_office_formats():
    doc = build_document(doc_type="docx", title="सूचना", paragraphs=["आज छुट्टी है"])
    from docx import Document

    text = "\n".join(p.text for p in Document(io.BytesIO(doc.content)).paragraphs)
    assert "आज छुट्टी है" in text


def test_an_excel_sheet_name_cannot_break_the_workbook():
    """Excel refuses to open a file whose sheet name holds []:*?/\\ - it does not
    warn, it just fails, so this is silently fatal if unhandled."""
    doc = build_document(doc_type="xlsx", title="Fees [2026]: Class 5/A?", rows=[["x"]])
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(doc.content))
    assert len(wb.active.title) <= 31
    assert not set("[]:*?/\\") & set(wb.active.title)
